"""
Общие фикстуры и генераторы образцов документов для тестов.

Фикстуры, требующие опциональных зависимостей, используют
``pytest.importorskip`` — тест пропускается, если библиотека не установлена
(соответствует философии «мягкой деградации» пакета).

Демо-данные едины (таблица + абзацы), латиницей — чтобы вывод не зависел от
наличия кириллических шрифтов в PDF-движке.
"""
from __future__ import annotations

import io
import json
import sqlite3
import zipfile

import pytest

from extractors import FileSource, build_default_extractor

# ── Единые демо-данные ─────────────────────────────────────────────────────
HEADING = "Quarterly Report"
PARA = "Sales by region for the first quarter."
COLS = ["Region", "Sales", "Growth"]
ROWS = [["North", "1200", "12"], ["South", "980", "-3"], ["West", "1560", "21"]]


@pytest.fixture
def svc():
    """Полностью собранный фасад (реестр + markitdown-рендер + OCR-заглушка)."""
    return build_default_extractor()


def source(data: bytes, filename: str) -> FileSource:
    """Короткий конструктор FileSource из байтов."""
    return FileSource(data=data, filename=filename)


# ── Генераторы образцов (байты), по одному формату ─────────────────────────
@pytest.fixture
def docx_bytes() -> bytes:
    docx = pytest.importorskip("docx")
    d = docx.Document()
    d.add_heading(HEADING, level=1)
    d.add_paragraph(PARA)
    t = d.add_table(rows=1, cols=3)
    for i, c in enumerate(COLS):
        t.rows[0].cells[i].text = c
    for row in ROWS:
        cells = t.add_row().cells
        for i, v in enumerate(row):
            cells[i].text = v
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


@pytest.fixture
def xlsx_bytes() -> bytes:
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(COLS)
    for row in ROWS:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = HEADING
    tbl = slide.shapes.add_table(len(ROWS) + 1, 3, Inches(0.5), Inches(1.8),
                                 Inches(8), Inches(3)).table
    for i, c in enumerate(COLS):
        tbl.cell(0, i).text = c
    for r, row in enumerate(ROWS, start=1):
        for i, v in enumerate(row):
            tbl.cell(r, i).text = v
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pdf_bytes() -> bytes:
    fitz = pytest.importorskip("fitz")
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), HEADING, fontsize=18)
    page.insert_text((72, 100), PARA, fontsize=11)
    y = 140
    for row in [COLS] + ROWS:
        page.insert_text((72, y), "  ".join(row), fontsize=11)
        y += 18
    data = doc.tobytes()
    doc.close()
    return data


@pytest.fixture
def pdf_blank_bytes() -> bytes:
    """Валидный PDF без текстового слоя (кандидат на OCR)."""
    fitz = pytest.importorskip("fitz")
    doc = fitz.open()
    doc.new_page()
    data = doc.tobytes()
    doc.close()
    return data


@pytest.fixture
def parquet_bytes() -> bytes:
    pd = pytest.importorskip("pandas")
    pytest.importorskip("pyarrow")
    df = pd.DataFrame(ROWS, columns=COLS)
    buf = io.BytesIO()
    df.to_parquet(buf, index=False)
    return buf.getvalue()


@pytest.fixture
def sqlite_path(tmp_path) -> str:
    p = tmp_path / "report.sqlite"
    con = sqlite3.connect(str(p))
    con.execute("CREATE TABLE sales(region TEXT, amount INTEGER, growth INTEGER)")
    con.executemany("INSERT INTO sales VALUES (?,?,?)",
                    [(r[0], int(r[1]), int(r[2])) for r in ROWS])
    con.execute("CREATE TABLE totals(metric TEXT, value INTEGER)")
    con.executemany("INSERT INTO totals VALUES (?,?)", [("total", 3740)])
    con.commit()
    con.close()
    return str(p)


@pytest.fixture
def csv_bytes() -> bytes:
    return ("\n".join([",".join(COLS)] + [",".join(r) for r in ROWS]) + "\n").encode()


@pytest.fixture
def tsv_bytes() -> bytes:
    return ("\n".join(["\t".join(COLS)] + ["\t".join(r) for r in ROWS]) + "\n").encode()


@pytest.fixture
def html_bytes() -> bytes:
    head = "<tr>" + "".join(f"<th>{c}</th>" for c in COLS) + "</tr>"
    body = "".join("<tr>" + "".join(f"<td>{v}</td>" for v in r) + "</tr>" for r in ROWS)
    return (f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{HEADING}</title></head>"
            f"<body><h1>{HEADING}</h1><p>{PARA}</p>"
            f"<table>{head}{body}</table></body></html>").encode()


@pytest.fixture
def json_bytes() -> bytes:
    obj = {"title": HEADING, "rows": [dict(zip(COLS, r)) for r in ROWS]}
    return json.dumps(obj, ensure_ascii=False).encode()


@pytest.fixture
def xml_bytes() -> bytes:
    rows = "".join("<row>" + "".join(f"<c>{v}</c>" for v in r) + "</row>" for r in ROWS)
    return (f"<?xml version='1.0' encoding='utf-8'?><report><title>{HEADING}</title>"
            f"<data>{rows}</data></report>").encode()


@pytest.fixture
def rtf_bytes() -> bytes:
    def rrow(cells):
        defs = "".join(f"\\cellx{3000 * (i + 1)}" for i in range(len(cells)))
        body = "".join(f"\\intbl {c}\\cell " for c in cells)
        return "\\trowd" + defs + body + "\\row\n"

    rtf = (r"{\rtf1\ansi\pard " + HEADING + r"\par " + PARA + r"\par " + "\n"
           + rrow(COLS) + "".join(rrow(r) for r in ROWS) + r"\pard End.\par }")
    return rtf.encode("ascii")


@pytest.fixture
def ipynb_bytes() -> bytes:
    nb = {
        "cells": [
            {"cell_type": "markdown", "metadata": {}, "source": [f"# {HEADING}\n", PARA]},
            {"cell_type": "code", "metadata": {}, "execution_count": 1, "outputs": [],
             "source": ["print('hi')\n"]},
        ],
        "metadata": {}, "nbformat": 4, "nbformat_minor": 5,
    }
    return json.dumps(nb).encode()


@pytest.fixture
def zip_bytes() -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("q1/sales.csv", "\n".join([",".join(COLS)] + [",".join(r) for r in ROWS]))
        z.writestr("notes/readme.md", f"# {HEADING}\n\n{PARA}\n")
    return buf.getvalue()


@pytest.fixture
def epub_bytes() -> bytes:
    head = "<tr>" + "".join(f"<th>{c}</th>" for c in COLS) + "</tr>"
    body = "".join("<tr>" + "".join(f"<td>{v}</td>" for v in r) + "</tr>" for r in ROWS)
    chapter = (f"<?xml version='1.0' encoding='utf-8'?>"
               f"<!DOCTYPE html><html xmlns='http://www.w3.org/1999/xhtml'><head>"
               f"<title>{HEADING}</title></head><body><h1>{HEADING}</h1><p>{PARA}</p>"
               f"<table>{head}{body}</table></body></html>")
    opf = (f"<?xml version='1.0'?><package xmlns='http://www.idpf.org/2007/opf' version='2.0' "
           f"unique-identifier='id'><metadata xmlns:dc='http://purl.org/dc/elements/1.1/'>"
           f"<dc:title>{HEADING}</dc:title><dc:identifier id='id'>urn:uuid:demo</dc:identifier>"
           f"<dc:language>en</dc:language></metadata>"
           f"<manifest><item id='ch1' href='chapter1.xhtml' media-type='application/xhtml+xml'/>"
           f"<item id='ncx' href='toc.ncx' media-type='application/x-dtbncx+xml'/></manifest>"
           f"<spine toc='ncx'><itemref idref='ch1'/></spine></package>")
    ncx = ("<?xml version='1.0'?><ncx xmlns='http://www.daisy.org/z3986/2005/ncx/' version='2005-1'>"
           "<head/><docTitle><text>R</text></docTitle><navMap><navPoint id='n1' playOrder='1'>"
           "<navLabel><text>Ch1</text></navLabel><content src='chapter1.xhtml'/></navPoint></navMap></ncx>")
    container = ("<?xml version='1.0'?><container version='1.0' "
                 "xmlns='urn:oasis:names:tc:opendocument:xmlns:container'><rootfiles>"
                 "<rootfile full-path='OEBPS/content.opf' media-type='application/oebps-package+xml'/>"
                 "</rootfiles></container>")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        z.writestr("META-INF/container.xml", container)
        z.writestr("OEBPS/content.opf", opf)
        z.writestr("OEBPS/toc.ncx", ncx)
        z.writestr("OEBPS/chapter1.xhtml", chapter)
    return buf.getvalue()
